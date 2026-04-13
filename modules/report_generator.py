import io
import datetime
import numpy as np
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.colors as mcolors
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


def _rgb(hex_color):
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


C_BG       = _rgb("0E1117")
C_CARD     = _rgb("1A1A2E")
C_ACCENT   = _rgb("4F8BF9")
C_WHITE    = _rgb("FAFAFA")
C_GRAY     = _rgb("8B8B9E")
C_GREEN    = _rgb("00C49A")
C_RED      = _rgb("FF6B6B")
C_BORDER   = _rgb("262730")


def _set_bg(slide, color=C_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(slide, text, left, top, width, height,
               font_size=14, bold=False, color=None,
               align=PP_ALIGN.LEFT, word_wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = word_wrap
    for i, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(font_size)
        r.font.bold = bold
        r.font.color.rgb = color or C_WHITE
    return txb


def _add_card(slide, val_text, label_text, left, top,
              width=Inches(2.9), height=Inches(1.25),
              val_color=C_ACCENT, border_color=C_ACCENT):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_CARD
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    _add_text(slide, val_text, left + Inches(0.1), top + Inches(0.08),
              width - Inches(0.2), Inches(0.65),
              font_size=26, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    _add_text(slide, label_text, left + Inches(0.1), top + Inches(0.72),
              width - Inches(0.2), Inches(0.45),
              font_size=11, color=C_GRAY, align=PP_ALIGN.CENTER)


def _add_divider(slide, left=Inches(0.5), top=Inches(0.95),
                 width=Inches(9.0), color=C_ACCENT):
    bar = slide.shapes.add_shape(1, left, top, width, Pt(2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def _add_accent_bar(slide):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.08), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = C_ACCENT
    bar.line.fill.background()


def _fig_to_stream(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    buf.seek(0)
    plt.close(fig)
    return buf


def _make_dist_chart(df, num_cols):
    cols = num_cols[:6]
    n = len(cols)
    ncols_g = min(3, n)
    nrows_g = (n + ncols_g - 1) // ncols_g
    fig, axes = plt.subplots(nrows_g, ncols_g, figsize=(12, 4.2 * nrows_g))
    fig.patch.set_facecolor("#0E1117")
    axes_flat = np.array(axes).flatten()
    for i, col in enumerate(cols):
        ax = axes_flat[i]
        ax.set_facecolor("#1A1A2E")
        data = df[col].dropna()
        ax.hist(data, bins=28, color="#4F8BF9", edgecolor="none", alpha=0.85)
        ax.axvline(data.mean(), color="#FF6B6B", lw=1.4, ls="--", alpha=0.85, label="mean")
        ax.set_title(col, color="white", fontsize=9, fontweight="bold", pad=5)
        ax.tick_params(colors="#8B8B9E", labelsize=7)
        for sp in ax.spines.values():
            sp.set_edgecolor("#262730")
    for j in range(i + 1, len(axes_flat)):
        axes_flat[j].set_visible(False)
    fig.tight_layout(pad=1.5)
    return _fig_to_stream(fig)


def _make_corr_chart(df, num_cols):
    corr = df[num_cols[:12]].corr()
    fig, ax = plt.subplots(figsize=(8, 5.5))
    fig.patch.set_facecolor("#0E1117")
    ax.set_facecolor("#0E1117")
    im = ax.imshow(corr.values, cmap="RdBu_r", vmin=-1, vmax=1, aspect="auto")
    ax.set_xticks(range(len(corr.columns)))
    ax.set_yticks(range(len(corr.columns)))
    ax.set_xticklabels(corr.columns, rotation=40, ha="right", color="white", fontsize=8)
    ax.set_yticklabels(corr.columns, color="white", fontsize=8)
    for i in range(len(corr)):
        for j in range(len(corr.columns)):
            v = corr.iloc[i, j]
            ax.text(j, i, f"{v:.2f}", ha="center", va="center",
                    color="white" if abs(v) > 0.3 else "#666688", fontsize=7)
    plt.colorbar(im, ax=ax, fraction=0.046, pad=0.04)
    fig.tight_layout()
    return _fig_to_stream(fig)


def _make_feat_imp_chart(feature_names, importances):
    idx = np.argsort(importances)[-min(12, len(importances)):]
    fig, ax = plt.subplots(figsize=(8, 3.8))
    fig.patch.set_facecolor("#0E1117")
    ax.set_facecolor("#1A1A2E")
    ax.barh([feature_names[i] for i in idx],
            [importances[i] for i in idx],
            color="#4F8BF9", alpha=0.85)
    ax.set_xlabel("Importance", color="white")
    ax.set_title("Feature Importance", color="white", fontweight="bold")
    ax.tick_params(colors="#8B8B9E")
    for sp in ax.spines.values():
        sp.set_edgecolor("#262730")
    fig.tight_layout()
    return _fig_to_stream(fig)


def generate_ppt(df, cleaned_df, insights, ml_results=None, file_name="dataset"):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    num_cols = df.select_dtypes(include=[np.number]).columns.tolist()
    cat_cols = df.select_dtypes(include=["object", "string", "category"]).columns.tolist()
    today = datetime.date.today().strftime("%B %d, %Y")

    # ── SLIDE 1: TITLE ────────────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    _add_accent_bar(s)

    _add_text(s, "AI Data Scientist Report",
              Inches(0.5), Inches(2.0), Inches(9.2), Inches(1.3),
              font_size=40, bold=True)
    _add_divider(s, top=Inches(3.4), width=Inches(9))
    _add_text(s, f"Dataset: {file_name}",
              Inches(0.5), Inches(3.65), Inches(9), Inches(0.6),
              font_size=22, color=C_ACCENT)
    _add_text(s,
              f"Generated: {today}  •  {df.shape[0]:,} rows × {df.shape[1]} columns",
              Inches(0.5), Inches(4.35), Inches(9), Inches(0.5),
              font_size=14, color=C_GRAY)
    _add_text(s, "Powered by Streamlit · Pandas · Plotly · scikit-learn",
              Inches(0.5), Inches(6.85), Inches(9), Inches(0.4),
              font_size=11, color=C_GRAY, align=PP_ALIGN.CENTER)

    # ── SLIDE 2: DATASET OVERVIEW ─────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    _add_text(s, "Dataset Overview", Inches(0.5), Inches(0.3), Inches(9), Inches(0.65),
              font_size=30, bold=True)
    _add_divider(s)

    cards = [
        ("Total Rows",          f"{df.shape[0]:,}",               C_ACCENT),
        ("Total Columns",       f"{df.shape[1]}",                 C_ACCENT),
        ("Numeric Columns",     f"{len(num_cols)}",               C_GREEN),
        ("Categorical Columns", f"{len(cat_cols)}",               C_GREEN),
        ("Missing Values",      f"{int(df.isnull().sum().sum()):,}", C_RED),
        ("Duplicate Rows",      f"{int(df.duplicated().sum()):,}", C_RED),
    ]
    for i, (label, val, color) in enumerate(cards):
        col = i % 3
        row = i // 3
        _add_card(s, val, label,
                  Inches(0.3 + col * 3.2), Inches(1.2 + row * 1.5),
                  val_color=color, border_color=color)

    col_lines = []
    for c in df.columns[:12]:
        col_lines.append(f"  • {c}  [{df[c].dtype}]  —  {df[c].nunique()} unique,  "
                         f"{int(df[c].isnull().sum())} null")
    _add_text(s, "Column Overview:", Inches(0.5), Inches(4.3), Inches(9.2), Inches(0.4),
              font_size=13, bold=True)
    _add_text(s, "\n".join(col_lines), Inches(0.5), Inches(4.7), Inches(9.2), Inches(2.4),
              font_size=10, color=C_GRAY)

    # ── SLIDE 3: DATA CLEANING ────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    _add_text(s, "Data Cleaning Summary", Inches(0.5), Inches(0.3), Inches(9), Inches(0.65),
              font_size=30, bold=True)
    _add_divider(s)

    bm  = int(df.isnull().sum().sum())
    am  = int(cleaned_df.isnull().sum().sum())
    bd  = int(df.duplicated().sum())
    ad  = int(cleaned_df.duplicated().sum())
    rem = df.shape[0] - cleaned_df.shape[0]
    fin = cleaned_df.shape[0]

    cleaning_cards = [
        ("Missing Before", f"{bm:,}",  C_RED),
        ("Missing After",  f"{am:,}",  C_GREEN),
        ("Dupes Before",   f"{bd:,}",  C_RED),
        ("Dupes After",    f"{ad:,}",  C_GREEN),
        ("Rows Removed",   f"{rem:,}", C_RED if rem > 0 else C_GRAY),
        ("Final Rows",     f"{fin:,}", C_GREEN),
    ]
    for i, (label, val, color) in enumerate(cleaning_cards):
        col = i % 3
        row = i // 3
        _add_card(s, val, label,
                  Inches(0.3 + col * 3.2), Inches(1.2 + row * 1.5),
                  val_color=color, border_color=color)

    notes = [
        "• Duplicate rows removed",
        "• Numeric NaNs filled with column median",
        "• Categorical NaNs filled with column mode",
        "• Dataset validated for completeness",
        f"• Final shape: {cleaned_df.shape[0]:,} rows × {cleaned_df.shape[1]} columns"
    ]
    _add_text(s, "\n".join(notes), Inches(0.5), Inches(4.4), Inches(9), Inches(2.4),
              font_size=13, color=C_GRAY)

    # ── SLIDE 4: DISTRIBUTIONS ────────────────────────────────────────────────
    if num_cols:
        s = prs.slides.add_slide(blank)
        _set_bg(s)
        _add_text(s, "Feature Distributions", Inches(0.5), Inches(0.3), Inches(9), Inches(0.65),
                  font_size=30, bold=True)
        _add_divider(s)
        buf = _make_dist_chart(cleaned_df, num_cols)
        s.shapes.add_picture(buf, Inches(0.3), Inches(1.1), Inches(9.4), Inches(5.9))

    # ── SLIDE 5: CORRELATION HEATMAP ──────────────────────────────────────────
    if len(num_cols) >= 2:
        s = prs.slides.add_slide(blank)
        _set_bg(s)
        _add_text(s, "Correlation Heatmap", Inches(0.5), Inches(0.3), Inches(9), Inches(0.65),
                  font_size=30, bold=True)
        _add_divider(s)
        buf = _make_corr_chart(cleaned_df, num_cols)
        s.shapes.add_picture(buf, Inches(0.5), Inches(1.1), Inches(9.0), Inches(5.9))

    # ── SLIDE 6: KEY INSIGHTS ─────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    _add_text(s, "Key Automated Insights", Inches(0.5), Inches(0.3), Inches(9), Inches(0.65),
              font_size=30, bold=True)
    _add_divider(s)

    if insights:
        mid = min(8, (len(insights) + 1) // 2)
        left_txt  = "\n".join(f"{i['icon']}  {i['message']}" for i in insights[:mid])
        right_txt = "\n".join(f"{i['icon']}  {i['message']}" for i in insights[mid:mid * 2])
        _add_text(s, left_txt,  Inches(0.5), Inches(1.1), Inches(4.5), Inches(5.9),
                  font_size=11, color=C_GRAY)
        if right_txt:
            _add_text(s, right_txt, Inches(5.1), Inches(1.1), Inches(4.5), Inches(5.9),
                      font_size=11, color=C_GRAY)
    else:
        _add_text(s, "✅  No major issues detected.\nThe dataset appears clean and well-structured.",
                  Inches(0.5), Inches(1.5), Inches(9), Inches(2),
                  font_size=16, color=C_GREEN)

    # ── SLIDE 7: ML RESULTS ───────────────────────────────────────────────────
    if ml_results:
        s = prs.slides.add_slide(blank)
        _set_bg(s)
        _add_text(s, "ML Prediction Results", Inches(0.5), Inches(0.3), Inches(9), Inches(0.65),
                  font_size=30, bold=True)
        _add_divider(s)

        task    = ml_results.get("task", "—")
        model   = ml_results.get("model", "—")
        target  = ml_results.get("target", "—")
        metrics = ml_results.get("metrics", {})
        feats   = ml_results.get("features", [])

        _add_text(s, f"Task: {task}   |   Model: {model}   |   Target: {target}",
                  Inches(0.5), Inches(1.1), Inches(9), Inches(0.45),
                  font_size=14, color=C_ACCENT)

        for i, (k, v) in enumerate(list(metrics.items())[:6]):
            col = i % 3
            row = i // 3
            v_str = f"{v:.4f}" if isinstance(v, float) else str(v)
            _add_card(s, v_str, k,
                      Inches(0.3 + col * 3.2), Inches(1.7 + row * 1.5),
                      val_color=C_GREEN, border_color=C_GREEN)

        if feats:
            feat_str = "Features: " + ", ".join(feats[:15])
            if len(feats) > 15:
                feat_str += f" (+{len(feats)-15} more)"
            _add_text(s, feat_str, Inches(0.5), Inches(5.45), Inches(9), Inches(0.45),
                      font_size=11, color=C_GRAY)

        imp = ml_results.get("feature_importances")
        feat_names = ml_results.get("feature_names", feats)
        if imp is not None and len(imp) == len(feat_names):
            buf = _make_feat_imp_chart(list(feat_names), list(imp))
            s.shapes.add_picture(buf, Inches(0.5), Inches(5.9), Inches(9.0), Inches(1.45))

    # ── SLIDE 8: CONCLUSION ───────────────────────────────────────────────────
    s = prs.slides.add_slide(blank)
    _set_bg(s)
    _add_accent_bar(s)

    _add_text(s, "Conclusion & Next Steps",
              Inches(0.5), Inches(0.35), Inches(9), Inches(0.75),
              font_size=30, bold=True)
    _add_divider(s)

    pts = [
        f"✅  {cleaned_df.shape[0]:,} clean rows × {cleaned_df.shape[1]} columns after preprocessing",
        f"📊  {len(num_cols)} numeric and {len(cat_cols)} categorical features analysed",
        f"💡  {len(insights)} automated insights discovered",
    ]
    if ml_results:
        metrics = ml_results.get("metrics", {})
        m_item = next(iter(metrics.items()), (None, None))
        if m_item[0]:
            pts.append(
                f"🤖  {ml_results.get('task', '')} model ({ml_results.get('model', '')}) — "
                f"{m_item[0]}: {m_item[1]:.4f}"
            )
    pts += [
        "🔍  Review outliers and skewed distributions for feature engineering",
        "📈  Tune hyperparameters and explore ensemble methods for higher accuracy",
        "🚀  Package the best model as a REST API for production deployment",
    ]
    for i, pt in enumerate(pts):
        _add_text(s, pt, Inches(0.7), Inches(1.25 + i * 0.72), Inches(8.8), Inches(0.65),
                  font_size=14, color=C_GRAY)

    _add_text(s, "Generated by AI Data Scientist Web App",
              Inches(0.5), Inches(6.9), Inches(9), Inches(0.4),
              font_size=10, color=_rgb("444460"), align=PP_ALIGN.CENTER)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
